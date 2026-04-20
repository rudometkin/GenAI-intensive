"""
Builds a .pptx file from a Presentation model using python-pptx.

Slide layout strategy:
- Slide 0 (title): large centered title + subtitle with generation date.
- Slides 1..N (content): heading + bullet list with theme-aware styling.
- Speaker notes are attached to every slide.

All dimensions use EMU (English Metric Units) via pptx.util helpers.
"""

import io
from datetime import date

from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from models import Presentation, ThemeConfig


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert '#1a365d' to an RGBColor object."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_title_slide(
    prs: PptxPresentation,
    title: str,
    subtitle: str,
    theme: ThemeConfig,
) -> None:
    """Create the opening title slide with centered text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(theme.primary_color)

    # Title text box
    left, top = Inches(0.8), Inches(2.0)
    width, height = Inches(8.4), Inches(1.8)
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = theme.font_heading
    p.alignment = PP_ALIGN.CENTER

    # Subtitle (date + tagline)
    sub_top = Inches(4.0)
    sub_box = slide.shapes.add_textbox(left, sub_top, width, Inches(0.8))
    sub_tf = sub_box.text_frame
    sub_tf.word_wrap = True

    sub_p = sub_tf.paragraphs[0]
    sub_p.text = f"{subtitle}  |  {date.today().strftime('%B %d, %Y')}"
    sub_p.font.size = Pt(18)
    sub_p.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    sub_p.font.name = theme.font_body
    sub_p.alignment = PP_ALIGN.CENTER

    # Speaker notes
    slide.notes_slide.notes_text_frame.text = subtitle


def _add_content_slide(
    prs: PptxPresentation,
    slide_num: int,
    title: str,
    bullets: list[str],
    notes: str,
    theme: ThemeConfig,
) -> None:
    """Create a content slide with heading and styled bullet list."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(theme.bg_color)

    # Accent bar on the left
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0), Inches(0),
        Inches(0.15), Inches(7.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex_to_rgb(theme.secondary_color)
    bar.line.fill.background()

    # Slide number badge
    num_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(0.6), Inches(0.5))
    num_tf = num_box.text_frame
    num_p = num_tf.paragraphs[0]
    num_p.text = str(slide_num)
    num_p.font.size = Pt(14)
    num_p.font.bold = True
    num_p.font.color.rgb = _hex_to_rgb(theme.secondary_color)
    num_p.font.name = theme.font_body

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(8.4), Inches(1.0))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True

    title_p = title_tf.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = _hex_to_rgb(theme.primary_color)
    title_p.font.name = theme.font_heading
    title_p.alignment = PP_ALIGN.LEFT

    # Bullets
    if bullets:
        bullet_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.6), Inches(8.0), Inches(4.8)
        )
        bullet_tf = bullet_box.text_frame
        bullet_tf.word_wrap = True

        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = bullet_tf.paragraphs[0]
            else:
                p = bullet_tf.add_paragraph()

            p.text = f"\u2022  {bullet_text}"
            p.font.size = Pt(18)
            p.font.color.rgb = _hex_to_rgb(theme.primary_color)
            p.font.name = theme.font_body
            p.space_after = Pt(12)
            p.alignment = PP_ALIGN.LEFT

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes


def build_pptx(presentation: Presentation, theme: ThemeConfig) -> bytes:
    """
    Convert a Presentation model into a downloadable .pptx byte stream.

    Args:
        presentation: Validated Presentation with slides.
        theme: ThemeConfig controlling colors and fonts.

    Returns:
        Bytes of the .pptx file (ready for st.download_button).
    """
    prs = PptxPresentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slides = presentation.slides
    if not slides:
        raise ValueError("Presentation has no slides to build.")

    # First slide is always the title slide
    first = slides[0]
    _add_title_slide(prs, presentation.title, first.notes or "", theme)

    # Content slides
    for idx, slide_content in enumerate(slides[1:], start=1):
        _add_content_slide(
            prs,
            slide_num=idx,
            title=slide_content.title,
            bullets=slide_content.bullets,
            notes=slide_content.notes,
            theme=theme,
        )

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
